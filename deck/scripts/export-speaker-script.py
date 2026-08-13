#!/usr/bin/env python3
"""Render speaker-script.pptx, speaker-script.docx, and a mapping README.

Reads the JSON payload written by export-speaker-script.mjs (spoken notes
already stripped of ⟢ caveats; silent slides already omitted).
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches, Pt as PptxPt

# Deck theme tokens from deck/src/styles/theme.css
BEIGE = PptxRGB(0xEF, 0xE6, 0xD6)
BLUE = PptxRGB(0x0D, 0x17, 0x34)
GOLD = PptxRGB(0xC8, 0x9B, 0x4B)
GOLD_DARK = PptxRGB(0x9A, 0x6C, 0x28)
INK = PptxRGB(0x0E, 0x15, 0x24)
BEIGE_DIM = PptxRGB(0xD8, 0xCB, 0xB4)
CREAM = PptxRGB(0xF5, 0xE8, 0xC8)

DOC_INK = RGBColor(0x0E, 0x15, 0x24)
DOC_BLUE = RGBColor(0x0D, 0x17, 0x34)
DOC_GOLD = RGBColor(0x9A, 0x6C, 0x28)
DOC_MUTED = RGBColor(0x6F, 0x75, 0x6F)

SLIDE_W = PptxInches(13.333)
SLIDE_H = PptxInches(7.5)


def _set_run_font(run, *, name: str, size_pt: float, color: PptxRGB, bold: bool = False) -> None:
    run.font.name = name
    run.font.size = PptxPt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    # Keep Google Slides from substituting a CJK face on Latin copy.
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn_a("latin"))
    if latin is None:
        latin = etree.SubElement(rPr, qn_a("latin"))
    latin.set("typeface", name)


def qn_a(tag: str) -> str:
    return f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}"


def _no_line(shape) -> None:
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, fill: PptxRGB):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    _no_line(shape)
    # Decorative bars should not contribute empty text frames when copied.
    sp = shape._element
    tx_body = sp.find("{http://schemas.openxmlformats.org/presentationml/2006/main}txBody")
    if tx_body is not None:
        sp.remove(tx_body)
    return shape


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tf


def body_size(word_count: int) -> float:
    if word_count > 70:
        return 18
    if word_count > 55:
        return 20
    if word_count > 42:
        return 22
    if word_count > 30:
        return 24
    return 26


def word_count(text: str) -> int:
    return len(text.split())


def build_pptx(payload: dict, dest: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = payload["spokenCount"]

    for slide_data in payload["slides"]:
        slide = prs.slides.add_slide(blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BEIGE

        _add_rect(slide, 0, 0, SLIDE_W, PptxInches(1.05), BLUE)
        _add_rect(slide, 0, PptxInches(1.05), SLIDE_W, PptxInches(0.06), GOLD)

        header = _textbox(slide, PptxInches(0.7), PptxInches(0.18), PptxInches(11.9), PptxInches(0.38))
        header.paragraphs[0].clear()
        p = header.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"SLIDE {slide_data['pptxIndex']}"
        _set_run_font(run, name="Calibri", size_pt=13, color=GOLD, bold=True)

        run = p.add_run()
        run.text = f"    {slide_data['id']}    ·    {slide_data['speaker']}"
        _set_run_font(run, name="Calibri", size_pt=13, color=CREAM, bold=False)

        title = _textbox(slide, PptxInches(0.7), PptxInches(0.52), PptxInches(11.9), PptxInches(0.44))
        tp = title.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        run = tp.add_run()
        run.text = slide_data["headline"]
        _set_run_font(run, name="Calibri", size_pt=16, color=BEIGE, bold=False)

        spoken = slide_data["spoken"]
        size = body_size(word_count(spoken))
        body = _textbox(slide, PptxInches(0.7), PptxInches(1.45), PptxInches(11.9), PptxInches(5.35))
        paragraphs = [part.strip() for part in spoken.split("\n") if part.strip()] or [spoken]
        body.paragraphs[0].clear()
        for index, text in enumerate(paragraphs):
            bp = body.paragraphs[0] if index == 0 else body.add_paragraph()
            bp.alignment = PP_ALIGN.LEFT
            bp.line_spacing = 1.32
            bp.space_after = PptxPt(10 if index < len(paragraphs) - 1 else 0)
            run = bp.add_run()
            run.text = text
            _set_run_font(run, name="Calibri", size_pt=size, color=INK, bold=False)

        notes = slide.notes_slide
        notes.notes_text_frame.text = spoken

        folio = _textbox(slide, PptxInches(11.4), PptxInches(6.95), PptxInches(1.3), PptxInches(0.32))
        fp = folio.paragraphs[0]
        fp.alignment = PP_ALIGN.RIGHT
        run = fp.add_run()
        run.text = f"{slide_data['pptxIndex']} / {total}"
        _set_run_font(run, name="Calibri", size_pt=11, color=GOLD_DARK, bold=False)

    prs.save(dest)


def _set_doc_run(run, *, name: str, size_pt: float, color: RGBColor, bold: bool = False) -> None:
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = etree.SubElement(rPr, qn("w:rFonts"))
    rFonts.set(qn("w:ascii"), name)
    rFonts.set(qn("w:hAnsi"), name)
    rFonts.set(qn("w:cs"), name)


def build_docx(payload: dict, dest: Path) -> None:
    doc = Document()
    section = doc.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)

    heading = doc.add_paragraph()
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    heading.paragraph_format.space_after = Pt(4)
    run = heading.add_run(payload["title"])
    _set_doc_run(run, name="Calibri", size_pt=22, color=DOC_BLUE, bold=True)

    kicker = doc.add_paragraph()
    kicker.paragraph_format.space_after = Pt(18)
    run = kicker.add_run(
        f"{payload['spokenCount']} spoken slides. "
        "Paste each block into Google Docs, or upload this file via File → Open."
    )
    _set_doc_run(run, name="Calibri", size_pt=11, color=DOC_MUTED)

    for slide_data in payload["slides"]:
        label = doc.add_paragraph()
        label.paragraph_format.space_before = Pt(18)
        label.paragraph_format.space_after = Pt(2)
        run = label.add_run(f"Slide {slide_data['pptxIndex']}")
        _set_doc_run(run, name="Calibri", size_pt=14, color=DOC_GOLD, bold=True)
        run = label.add_run(f"  ·  {slide_data['id']}")
        _set_doc_run(run, name="Calibri", size_pt=12, color=DOC_MUTED)

        title = doc.add_paragraph()
        title.paragraph_format.space_after = Pt(8)
        run = title.add_run(slide_data["headline"])
        _set_doc_run(run, name="Calibri", size_pt=16, color=DOC_BLUE, bold=True)

        body = doc.add_paragraph()
        body.paragraph_format.space_after = Pt(12)
        body.paragraph_format.line_spacing = 1.35
        run = body.add_run(slide_data["spoken"])
        _set_doc_run(run, name="Calibri", size_pt=14, color=DOC_INK)

        spacer = doc.add_paragraph()
        spacer.paragraph_format.space_after = Pt(6)

    doc.save(dest)


def build_readme(payload: dict, dest: Path) -> None:
    lines = [
        "# Speaker script mapping",
        "",
        "Spoken copy from `deck/src/slides/index.ts` `notes`, processed like `spokenNotes()`:",
        "everything after `⟢` is stripped. Slides whose spoken remainder is empty are omitted.",
        "PPTX/DOC numbering is the remapped spoken order (PPTX 1 = first slide they actually say).",
        "`original index` is the 1-based position in the deck registry (same number the",
        f"presenter overlay shows as `N / {payload.get('deckSlides', '?')}`).",
        "",
    ]
    for slide in payload["slides"]:
        lines.append(
            f"PPTX {slide['pptxIndex']} = deck id {slide['id']} "
            f"(original index {slide['originalIndex']}) — {slide['headline']}"
        )
    lines += [
        "",
        "## Import",
        "",
        "- **Google Slides:** File → Open → upload `speaker-script.pptx`, or File → Import slides.",
        "- **Google Docs:** File → Open → upload `speaker-script.docx`. Each slide is a heading plus the spoken paragraph.",
        "",
        "Regenerate with `node scripts/export-speaker-script.mjs` from `deck/`.",
        "",
    ]
    dest.write_text("\n".join(lines), encoding="utf-8")


def main() -> int:
    json_path = Path(sys.argv[1])
    out_dir = Path(sys.argv[2])
    payload = json.loads(json_path.read_text(encoding="utf-8"))
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = out_dir / "speaker-script.pptx"
    docx_path = out_dir / "speaker-script.docx"
    readme_path = out_dir / "README.md"
    build_pptx(payload, pptx_path)
    build_docx(payload, docx_path)
    build_readme(payload, readme_path)
    print(f"Wrote {pptx_path}")
    print(f"Wrote {docx_path}")
    print(f"Wrote {readme_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
